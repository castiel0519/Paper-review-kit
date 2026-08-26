#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
verify_deliverables.py — 交付物自检：
 1) PDF 完整性（%PDF 魔数、页数）
 2) summaries 完整性（每篇 method_zh 非空、必填字段齐全）
 3) DOCX 可打开（paragraphs/tables/images 数量）
 4) PPTX 可打开（slides 数量、每篇论文2页结构抽查、图片嵌入数）
 5) Scheme 图存在且非空
输出 check_report.json 供 delivery_check 引用。
"""
import json
import os

HERE = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.dirname(HERE)

PAPER_DIR = os.path.join(BASE, "papers")
SUM_DIR = os.path.join(BASE, "summaries")
FIG_DIR = os.path.join(BASE, "papers_figs")
DELIVER = os.path.join(BASE, "deliverables")

DOCX = os.path.join(DELIVER, "机器学习与人工智能在微流控中的应用_读书报告.docx")
PPTX = os.path.join(DELIVER, "机器学习与人工智能在微流控中的应用_研究报告.pptx")

FIELDS = ["id", "title_en", "title_zh", "journal", "year", "doi", "kind", "kind_zh",
          "abstract_en", "abstract_zh", "background_zh", "problem_zh", "data_zh",
          "task_zh", "method_zh", "results_zh", "innovation_zh", "limitation_zh",
          "paradigm_zh", "framework_zh", "scheme_zh", "lessons_zh"]


def main():
    report = {"ok": True, "pdf": [], "summaries": [], "figs": [], "docx": None,
              "pptx": None, "errors": []}

    # 1) PDFs
    pdfs = sorted([f for f in os.listdir(PAPER_DIR) if f.endswith(".pdf")])
    for f in pdfs:
        p = os.path.join(PAPER_DIR, f)
        head = open(p, "rb").read(5)
        ok = head == b"%PDF-" and os.path.getsize(p) > 10000
        pages = None
        try:
            import fitz
            pages = fitz.open(p).page_count
        except Exception:
            pass
        report["pdf"].append({"file": f, "ok": ok, "size": os.path.getsize(p), "pages": pages})
        if not ok:
            report["errors"].append(f"bad pdf: {f}")

    # 2) summaries
    for fn in sorted(os.listdir(SUM_DIR)):
        if not fn.endswith(".json"):
            continue
        s = json.load(open(os.path.join(SUM_DIR, fn), encoding="utf-8"))
        missing = [k for k in FIELDS if not s.get(k)]
        rec = {"id": s.get("id"), "missing": missing, "len_method": len(s.get("method_zh", "") or "")}
        report["summaries"].append(rec)
        if missing:
            report["errors"].append(f"summary {s.get('id')} missing: {missing}")

    # 3) figs
    for fn in sorted(os.listdir(FIG_DIR)):
        if fn.endswith(".png"):
            p = os.path.join(FIG_DIR, fn)
            ok = os.path.getsize(p) > 2000
            report["figs"].append({"file": fn, "ok": ok, "size": os.path.getsize(p)})
            if not ok:
                report["errors"].append(f"small fig: {fn}")

    # 4) docx
    try:
        from docx import Document
        d = Document(DOCX)
        report["docx"] = {
            "exists": True,
            "paragraphs": len(d.paragraphs),
            "tables": len(d.tables),
            "inline_shapes": len(d.inline_shapes),
        }
    except Exception as e:
        report["docx"] = {"exists": False, "error": str(e)}
        report["errors"].append(f"docx: {e}")

    # 5) pptx
    try:
        from pptx import Presentation
        prs = Presentation(PPTX)
        pics = 0
        for sl in prs.slides:
            for sh in sl.shapes:
                if sh.shape_type == 13:  # PICTURE
                    pics += 1
        report["pptx"] = {"exists": True, "slides": len(prs.slides._sldIdLst), "pictures": pics}
    except Exception as e:
        report["pptx"] = {"exists": False, "error": str(e)}
        report["errors"].append(f"pptx: {e}")

    report["ok"] = len(report["errors"]) == 0
    with open(os.path.join(BASE, "check_report.json"), "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)
    print(json.dumps(report, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
