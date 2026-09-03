#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
make_pptx.py — 依据 summaries/summary_*.json 生成详细研究报告 PPT（16:9）。
结构（每个论文2页）：
  Page A：论文信息 + 内容讲解（背景/数据/方法/结果/创新点） + 论文 Scheme 图
  Page B：研究范式（范式标签 + 解析） + 研究框架解析（数据流 → 模型 → 训练 → 评估）
另加：封面 / 目录 / 背景与方法 / 领域全景 / 横向对比 / 趋势 / 结论 / 参考文献。
标题、年份范围、文件名、叙事文案由 project.yml 驱动。
输出：deliverables/{title}_研究报告.pptx
"""
import math
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from prk_config import (
    cfg_get, load_papers_meta, output_dir, parse_project_arg, project_range,
    project_title, read_json,
)
from prk_schema import validate_summary

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2E, 0x74, 0xB5)
LIGHT = RGBColor(0xEA, 0xF0, 0xF7)
GRAY = RGBColor(0x59, 0x59, 0x59)
DARK = RGBColor(0x26, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC5, 0x8F, 0x1B)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)

FONT = "微软雅黑"
SW, SH = 13.333, 7.5

try:
    from PIL import Image as PILImage
    HAS_PIL = True
except Exception:
    HAS_PIL = False


def set_font(run, size=11, bold=False, color=DARK, name=FONT, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def add_box(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.0)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, items, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, space_after=4):
    """items: list of (text, size, bold, color) 或 (text, size, bold)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, it in enumerate(items):
        text, size, bold = it[0], it[1], it[2]
        color = it[3] if len(it) > 3 else DARK
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = text
        set_font(r, size=size, bold=bold, color=color)
    return tb


def est_lines(text, width_in, font_pt):
    cpl = max(1, int(width_in * 72 / (font_pt * 1.02)))
    n = 0
    for seg in str(text).split("\n"):
        n += max(1, math.ceil(len(seg) / cpl))
    return n


def fit_size(text, width_in, height_in, base=11.5, floor=8.0):
    s = base
    while s > floor and est_lines(text, width_in, s) * (s * 1.42) / 72 > height_in:
        s -= 0.5
    return max(s, floor)


def add_bullets(slide, x, y, w, h, bullets, base=11.5, lead_color=BLUE, body_color=DARK):
    """bullets: list of (lead, body) —— lead 加粗，body 跟随"""
    size = fit_size("".join([b[0] + b[1] for b in bullets]), w, h, base=base)
    items = []
    for lead, body in bullets:
        items.append((f"▪ {lead}   ", size, True, lead_color))
        # body 与 lead 在同一段落较长时拆段
        if body:
            items.append((body, size, False, body_color))
    add_text(slide, x, y, w, h, items, line_spacing=0.92, space_after=3)


def add_title_bar(slide, title, subtitle=None):
    add_box(slide, 0, 0, SW, 0.95, fill=NAVY)
    add_text(slide, 0.45, 0.10, 10.6, 0.78, [(title, 21, True, WHITE)],
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, 11.0, 0.10, 2.05, 0.78, [(subtitle, 11, True, RGBColor(0xBF, 0xD7, 0xEE))],
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_pic_fit(slide, path, x, y, w, h):
    """在 box 内按比例放入图片，返回 (ix, iy, iw, ih)"""
    if not (path and os.path.exists(path)):
        add_box(slide, x, y, w, h, fill=LIGHT, line=BLUE)
        add_text(slide, x, y, w, h, [("Scheme 图暂缺（非OA或未下载PDF）", 11, True, GRAY)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return None
    try:
        if HAS_PIL:
            iw, ih = PILImage.open(path).size
        else:
            iw, ih = 4, 3
    except Exception:
        iw, ih = 4, 3
    scale = min(w / iw, h / ih)
    nw, nh = iw * scale * 96, ih * scale * 96  # 96dpi -> inches
    ix = x + (w - nw / 96) / 2
    iy = y + (h - nh / 96) / 2
    slide.shapes.add_picture(str(path), Inches(ix), Inches(iy), Inches(nw / 96), Inches(nh / 96))
    return (ix, iy, nw / 96, nh / 96)


def _report_lines(cfg, key, default):
    val = cfg_get(cfg, "report", "narrative", key)
    if isinstance(val, list) and val:
        return [str(x) for x in val]
    if isinstance(val, str) and val.strip():
        return [val]
    return default


def _generic_trends(cfg):
    return _report_lines(cfg, "trends", [
        "数据驱动建模：以可复用数据集与基准测试为底座，形成可比较、可复现的实验体系。",
        "端到端深度学习：以图像、序列或多模态信号为直接输入，完成检测、分类、分割与预测。",
        "物理信息与控制驱动学习：将领域先验、仿真器或闭环控制嵌入学习过程。",
        "生成式与自主实验：由目标反向设计结构与条件，LLM/智能体逐步参与实验设计。",
        "多模态集成：融合多源信息，推动从研究原型到应用落地。",
    ])


def _generic_conclusions(cfg):
    return _report_lines(cfg, "conclusions", [
        "AI 已成为本领域方法创新的核心引擎，覆盖设计、识别、控制与预测全链条。",
        "端到端深度学习与自主控制并行发展，二者互补形成闭环。",
        "多模态与高吞吐方向成为主流，现场/即时检测场景加速落地。",
        "工程落地需要可解释、可重复与低算力，开放数据与标准接口是关键。",
        "综述与基础模型推动「数据-算法-硬件-应用」方法学范式化。",
    ])


def _kind_map(cfg, summaries):
    """kind -> (zh, en, color, group)；配置优先，meta 派生兜底。"""
    colors = [BLUE, GREEN, GOLD, RGBColor(0x7B, 0x5E, 0xB7), RGBColor(0x1B, 0x6C, 0x8F),
              NAVY, RGBColor(0x9C, 0x50, 0x26), GREEN]
    configured = cfg_get(cfg, "report", "kinds", default=None)
    out = {}
    idx = 0
    for s in summaries:
        kind = str(s.get("kind") or "")
        if not kind or kind in out:
            continue
        zh = str(s.get("kind_zh") or kind)
        if isinstance(configured, dict) and kind in configured:
            c = configured[kind]
            en = str(c.get("en", ""))
            group = str(c.get("group", zh))
            color = c.get("color", "BLUE")
            color = {"BLUE": BLUE, "GREEN": GREEN, "GOLD": GOLD, "NAVY": NAVY}.get(color, BLUE)
        else:
            en = zh.replace("与", " & ").replace("（", "(").replace("）", ")") or "General"
            group = zh
            color = colors[idx % len(colors)]
        out[kind] = {"zh": zh, "en": en, "group": group, "color": color}
        idx += 1
    return out


def _kind_name(kind_map, kind, lang="zh"):
    info = kind_map.get(str(kind)) or {}
    return info.get("zh" if lang == "zh" else "en", str(kind or ""))


def slide_cover(prs, title, topic_en, rng, count, kind_text):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(s, 0, 0, SW, SH, fill=NAVY)
    add_box(s, 0.55, 1.35, 3.1, 0.14, fill=GOLD)
    add_text(s, 0.55, 1.75, 11.5, 1.5, [(title, 36, True, WHITE)], space_after=0)
    add_text(s, 0.55, 3.15, 11.5, 0.5, [(topic_en, 17, False, RGBColor(0xBF, 0xD7, 0xEE))])
    add_text(s, 0.55, 3.85, 11.5, 0.9, [
        (f"SCI 论文精读调研报告 ｜ {rng} ｜ 共 {count} 篇 ｜ 中文为主", 15, True, GOLD),
        (f"覆盖：{kind_text}", 11.5, False, RGBColor(0xBF, 0xD7, 0xEE)),
    ], space_after=8)
    add_text(s, 0.55, 6.6, 11.5, 0.5, [("每篇论文 2 页：内容讲解+Scheme ／ 研究范式+框架解析", 12, False, WHITE)])
    return s


def slide_toc(prs, count):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "目录  Contents")
    items = [
        ("01", "调研背景与选文标准", "SCI 论文调研方法"),
        ("02", "领域全景", "按研究方向分组"),
        ("03", "论文精读（每篇2页）", "内容讲解 + Scheme ／ 研究范式 + 框架解析"),
        ("04", "横向对比", "数据规模 · 任务 · 方法 · 核心指标"),
        ("05", "范式总结与趋势", "从传统模型到端到端学习与自主实验"),
        ("06", "结论与展望", "研究落地的关键路径"),
        ("07", "参考文献", f"{count} 篇论文完整题录"),
    ]
    y = 1.35
    for num, t1, t2 in items:
        add_box(s, 0.8, y, 0.62, 0.62, fill=BLUE)
        add_text(s, 0.8, y, 0.62, 0.62, [(num, 15, True, WHITE)], align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.6, y - 0.03, 10.6, 0.75, [
            (t1, 15, True, NAVY),
            (t2, 10.5, False, GRAY),
        ], space_after=1)
        y += 0.82
    return s


def slide_bg(prs, cfg, topic, rng):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "01 调研背景与选文标准")
    lines = _report_lines(cfg, "background", [
        f"{topic}领域的数据与实验空间庞大、反馈链路长，AI/ML 正成为加速设计、提升自动化与智能决策的关键工具。",
        f"范围：聚焦 {rng} 年 SCI 期刊论文，覆盖多个代表性研究方向。",
        "方法：OpenAlex / Crossref / Europe PMC / Unpaywall / 出版社官网检索核验元数据；优先合规来源获取 PDF；按统一模板提取结构化信息。",
        "精读维度：摘要 / 背景与问题 / 数据与任务 / 方法 / 关键结果 / 创新与局限 / 研究范式 / 研究框架解析 / Scheme解读 / 启示。",
    ])
    add_bullets(s, 0.7, 1.25, 12.0, 5.6, [("", x) for x in lines], base=12)
    return s


def slide_panorama(prs, summaries, kind_map):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    groups = [(k, v["group"], v["color"]) for k, v in kind_map.items()]
    if not groups:
        groups = [("general", "论文调研", BLUE)]
    add_title_bar(s, "02 领域全景：研究方向分布")
    w, h = 3.05, 2.35
    pos = [(0.45, 1.35), (3.68, 1.35), (6.91, 1.35), (10.14, 1.35),
           (0.45, 4.05), (3.68, 4.05), (6.91, 4.05), (10.14, 4.05)]
    for (kind, group, col), (x, y) in zip(groups[:8], pos[:len(groups)]):
        add_box(s, x, y, w, h, fill=LIGHT, line=col)
        add_box(s, x, y, w, 0.5, fill=col)
        add_text(s, x, y, w, 0.5, [(group[:14], 12.5, True, WHITE)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cnt = sum(1 for sm in summaries if str(sm.get("kind")) == kind)
        zh = kind_map.get(kind, {}).get("zh", group)
        add_text(s, x + 0.12, y + 0.6, w - 0.24, 1.6, [
            (f"代表论文 {cnt} 篇", 11, True, col),
            (zh[:30], 9.5, False, DARK),
        ], space_after=5)
    return s


def _resolve_scheme(s, fig_dir, cfg=None):
    """优先使用 summary.scheme_image，缺失时回退到 fig1。"""
    raw = s.get("scheme_image") or ""
    if raw:
        p = Path(raw)
        if not p.is_absolute():
            # 兼容两种常见写法：papers_figs/xx.png 或 xx.png
            try:
                if cfg:
                    root_rel = Path(cfg.get("_project_root", ".")) / p
                    if root_rel.exists():
                        return root_rel
            except Exception:
                pass
            p = fig_dir / p
        if p.exists():
            return p
    default = fig_dir / f"{s.get('id', '')}_fig1.png"
    return default if default.exists() else None


def slide_paper_a(prs, s, fig_dir, kind_map, cfg=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kind_en = _kind_name(kind_map, s.get("kind"), "en")
    add_title_bar(slide, f"论文 {s['id']} ｜ {str(s.get('title_zh', ''))[:26]}",
                  f"{kind_en} · {s.get('kind_zh', '')}")
    add_text(slide, 0.55, 1.02, 12.2, 0.4, [
        (f"{s.get('journal', '')} · {s.get('year', '')} ｜ DOI: {s.get('doi', '')} ｜ 英文题名：{str(s.get('title_en', ''))[:70]}",
         9.5, False, GRAY)])
    add_bullets(slide, 0.5, 1.5, 6.55, 5.4, [
        ("背景与问题：", str(s.get("background_zh", "")) + (" " + str(s.get("problem_zh", "")) if s.get("problem_zh") else "")),
        ("数据与任务：", str(s.get("data_zh", "")) + (" " + str(s.get("task_zh", "")) if s.get("task_zh") else "")),
        ("方法要点：", str(s.get("method_zh", ""))),
        ("关键结果：", str(s.get("results_zh", ""))),
        ("创新点：", str(s.get("innovation_zh", ""))),
    ], base=11.5)
    y = 6.92
    ms = (s.get("metrics_zh") or [])[:3]
    x = 0.5
    for m in ms:
        label = m.get("label", "") if isinstance(m, dict) else str(m)
        value = m.get("value", "") if isinstance(m, dict) else ""
        w = 2.1
        add_box(slide, x, y, w, 0.5, fill=LIGHT, line=BLUE)
        add_text(slide, x, y, w, 0.5, [(f"{label}  {value}", 9.5, True, NAVY)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += w + 0.12
    add_text(slide, 7.25, 1.5, 5.6, 0.35, [("论文 Scheme", 13, True, NAVY)])
    add_pic_fit(slide, _resolve_scheme(s, fig_dir, cfg), 7.25, 1.88, 5.6, 3.9)
    add_text(slide, 7.25, 5.85, 5.6, 1.5, [
        ("Scheme 解读：" + (str(s.get("scheme_zh", "")) or "见报告正文"), 10, False, DARK),
    ], space_after=2)
    return slide


def chip_row(slide, y, tags, x0=0.55, w=2.45):
    x = x0
    for t in tags[:5]:
        bw = max(1.0, min(2.4, 0.3 + 0.24 * len(t)))
        add_box(slide, x, y, bw, 0.45, fill=BLUE if t else LIGHT)
        add_text(slide, x, y, bw, 0.45, [(t, 10, True, WHITE)], align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += bw + 0.15
        if x > 12.8:
            break


def slide_paper_b(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, f"论文 {s['id']} ｜ {str(s.get('title_zh', ''))[:22]}",
                  "研究范式 Research Paradigm")
    tags = s.get("paradigm_tags") or []
    add_text(slide, 0.55, 1.05, 3.0, 0.4, [("范式标签", 11, True, NAVY)])
    chip_row(slide, 1.5, tags)
    add_box(slide, 0.55, 2.12, 12.2, 1.15, fill=LIGHT, line=BLUE)
    add_text(slide, 0.75, 2.22, 11.8, 0.35, [("研究范式解析", 12, True, NAVY)])
    add_text(slide, 0.75, 2.58, 11.8, 0.65, [(str(s.get("paradigm_zh", "")), 10.5, False, DARK)])
    add_text(slide, 0.55, 3.42, 6.0, 0.35, [("研究框架解析：数据流 → 模型 → 训练 → 评估 / 解读", 12, True, NAVY)])
    steps = s.get("framework_steps") or ["数据与预处理", "模型设计", "训练策略", "评估与解读"]
    bw, gap = 2.75, 0.35
    x0 = 0.55
    for i, st in enumerate(steps):
        x = x0 + i * (bw + gap)
        add_box(slide, x, 3.82, bw, 1.7, fill=LIGHT, line=BLUE)
        head = str(st)[:14]
        body = str(st)[14:] if len(str(st)) > 14 else ""
        add_text(slide, x + 0.08, 3.92, bw - 0.16, 1.5, [
            (f"STEP {i+1} {head}", 10.5, True, NAVY),
            (body, 9.5, False, DARK),
        ], space_after=3)
        if i < len(steps) - 1:
            add_text(slide, x + bw + 0.02, 4.35, gap + 0.05, 0.5, [("▶", 12, True, BLUE)],
                     align=PP_ALIGN.CENTER)
    add_box(slide, 0.55, 5.72, 6.0, 1.5, fill=RGBColor(0xEF, 0xF5, 0xEA), line=GREEN)
    add_text(slide, 0.72, 5.82, 5.6, 0.3, [("框架总评", 11.5, True, GREEN)])
    add_text(slide, 0.72, 6.12, 5.66, 1.0, [(str(s.get("framework_zh", "")), 9.5, False, DARK)])
    add_box(slide, 6.75, 5.72, 6.0, 1.5, fill=RGBColor(0xFB, 0xF1, 0xE2), line=GOLD)
    add_text(slide, 6.92, 5.82, 5.6, 0.3, [("对本领域研究的启示", 11.5, True, GOLD)])
    add_text(slide, 6.92, 6.12, 5.66, 1.0, [(str(s.get("lessons_zh", "")), 9.5, False, DARK)])
    return slide


def slide_compare(prs, summaries):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "04 横向对比：任务 · 数据 · 方法 · 核心指标")
    rows, cols = len(summaries) + 1, 5
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.1),
                                       Inches(12.5), Inches(6.0))
    tbl = tbl_shape.table
    widths = [1.6, 1.9, 2.8, 3.2, 3.0]
    for j, w in enumerate(widths):
        tbl.columns[j].width = Inches(w)
    heads = ["论文", "方向", "数据与任务", "方法", "核心指标"]
    for j, h in enumerate(heads):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                set_font(r, size=10, bold=True, color=WHITE)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
    tbl.rows[0].height = Inches(0.32)
    for i, s in enumerate(summaries, start=1):
        m0 = (s.get("metrics_zh") or [{}])[0]
        label = m0.get("label", "") if isinstance(m0, dict) else str(m0)
        value = m0.get("value", "") if isinstance(m0, dict) else ""
        vals = [
            f"{s['id']} {str(s.get('title_zh',''))[:14]}",
            str(s.get("kind_zh", ""))[:12],
            (str(s.get("data_zh", "")) + " / " + str(s.get("task_zh", "")))[:55],
            str(s.get("method_zh", ""))[:72],
            f"{label}: {value}"[:30],
        ]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.text = str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    set_font(r, size=7.5, bold=(j == 0), color=DARK)
            if i % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = LIGHT
        tbl.rows[i].height = Inches(0.36)
    return slide


def slide_trends(prs, lines):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "05 研究范式总结与发展趋势")
    y = 1.25
    for i, line in enumerate(lines):
        col = BLUE if i % 2 == 0 else GREEN
        add_box(s, 0.55, y, 12.2, 1.05, fill=LIGHT, line=col)
        add_text(s, 0.75, y + 0.06, 11.6, 0.9, [(line, 11.5, False, DARK)],
                 anchor=MSO_ANCHOR.MIDDLE)
        y += 1.14
        if i >= 4:
            break
    return s


def slide_conclusion(prs, lines):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "06 结论与展望")
    add_bullets(s, 0.7, 1.3, 12.0, 5.5, [("", x) for x in lines[:5]], base=13, lead_color=NAVY)
    return s


def slide_refs(prs, summaries, count):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, f"07 参考文献（{count}篇）")
    half = (len(summaries) + 1) // 2
    for i, sm in enumerate(summaries):
        col = 0 if i < half else 1
        x = 0.55 + col * 6.3
        yy = 1.15 + (i % half) * 0.62
        add_text(s, x, yy, 6.0, 0.6, [
            (f"[{i+1}] {str(sm.get('title_en',''))[:60]}…  {sm.get('journal','')} {sm.get('year','')}",
             7.5, False, DARK)], space_after=0)
    return s


def main():
    cfg, args = parse_project_arg()
    meta = load_papers_meta(cfg)
    title = project_title(cfg, meta)
    rng = project_range(cfg, meta)
    topic = cfg_get(cfg, "project", "topic", default="论文调研")
    topic_en = cfg_get(cfg, "project", "topic_en", default="Paper Review")

    sum_dir = output_dir(cfg, "summaries")
    fig_dir = output_dir(cfg, "papers_figs")
    del_dir = output_dir(cfg, "deliverables")

    summaries = []
    if sum_dir.is_dir():
        for fn in sorted(sum_dir.iterdir()):
            if fn.suffix != ".json":
                continue
            s = read_json(fn, default=None)
            if not isinstance(s, dict) or not s.get("id"):
                continue
            errors, _ = validate_summary(s)
            if errors:
                print(f"  [warn] {fn.name}: " + "; ".join(errors))
            summaries.append(s)
    summaries.sort(key=lambda s: str(s["id"]))

    kind_map = _kind_map(cfg, summaries)
    kind_text = " / ".join(v["group"] for v in kind_map.values()) or topic

    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    slide_cover(prs, title, topic_en, rng, len(summaries), kind_text)
    slide_toc(prs, len(summaries))
    slide_bg(prs, cfg, topic, rng)
    slide_panorama(prs, summaries, kind_map)
    for s in summaries:
        slide_paper_a(prs, s, fig_dir, kind_map, cfg)
        slide_paper_b(prs, s)
    slide_compare(prs, summaries)
    slide_trends(prs, _generic_trends(cfg))
    slide_conclusion(prs, _generic_conclusions(cfg))
    slide_refs(prs, summaries, len(summaries))
    out = del_dir / f"{title}_研究报告.pptx"
    prs.save(str(out))
    print("PPTX written:", out, "slides:", len(prs.slides._sldIdLst), "papers:", len(summaries))


if __name__ == "__main__":
    main()

