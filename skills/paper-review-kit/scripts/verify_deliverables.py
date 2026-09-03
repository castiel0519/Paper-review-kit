#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
verify_deliverables.py — 交付物自检：
 1) PDF 完整性（%PDF 魔数、页数）
 2) summaries 完整性（必填字段、pdf_status 枚举、长度告警）
 3) DOCX 可打开（paragraphs/tables/images 数量）
 4) PPTX 可打开（slides 数量、每篇论文2页结构、图片嵌入数）
 5) Scheme 图存在且非空
输出 check_report.json 供 delivery_check 引用。
文件名与路径由 project.yml / papers_meta 驱动。
"""
import os
from pathlib import Path

from prk_config import load_papers_meta, output_dir, parse_project_arg, project_path, project_title, read_json, write_json
from prk_schema import PDF_STATUSES, validate_summaries_dir

REQUIRED_SLIDE_COUNT = 8  # 封面/目录/背景/全景/对比/趋势/结论/文献


def pdf_check(path):
    try:
        size = os.path.getsize(path)
        with open(path, "rb") as f:
            head = f.read(5)
        pages = None
        try:
            import fitz
            doc = fitz.open(str(path))
            pages = doc.page_count
            doc.close()
        except Exception:
            pass
        return {
            "file": os.path.basename(path), "ok": head == b"%PDF-" and size > 10_000,
            "size": size, "pages": pages,
        }
    except Exception as e:
        return {"file": os.path.basename(path), "ok": False, "error": str(e)}


def main():
    cfg, args = parse_project_arg()
    meta = load_papers_meta(cfg)
    title = project_title(cfg, meta)

    paper_dir = output_dir(cfg, "papers")
    sum_dir = output_dir(cfg, "summaries")
    fig_dir = output_dir(cfg, "papers_figs")
    del_dir = output_dir(cfg, "deliverables")
    docx = del_dir / f"{title}_读书报告.docx"
    pptx = del_dir / f"{title}_研究报告.pptx"

    report = {"ok": True, "title": title, "pdf": [], "summaries": [],
              "figs": [], "docx": None, "pptx": None, "errors": []}

    # 1) PDFs
    for f in sorted(p for p in paper_dir.iterdir() if p.suffix.lower() == ".pdf"):
        rec = pdf_check(f)
        report["pdf"].append(rec)
        if not rec.get("ok"):
            report["errors"].append(f"bad pdf: {rec['file']}")

    # 2) summaries
    results = validate_summaries_dir(sum_dir)
    for fn in sorted(results):
        errors, warnings = results[fn]
        rec = {"file": fn, "errors": errors, "warnings": warnings}
        report["summaries"].append(rec)
        if errors:
            report["errors"].append(f"summary {fn}: " + "; ".join(errors))

    # 3) figs
    for f in sorted(p for p in fig_dir.iterdir() if p.suffix.lower() == ".png"):
        size = f.stat().st_size
        ok = size > 2000
        report["figs"].append({"file": f.name, "ok": ok, "size": size})
        if not ok:
            report["errors"].append(f"small fig: {f.name}")

    # 3.5) scheme 图存在性与视觉确认
    figs_info = read_json(project_path(cfg, "papers_figs", "figs_info.json"), default={})
    summaries_paths = sorted(p for p in sum_dir.iterdir() if p.suffix == ".json")
    for path in summaries_paths:
        s = read_json(path, default={})
        if not isinstance(s, dict):
            continue
        pid = str(s.get("id", ""))
        chosen = s.get("scheme_image") or ""
        if chosen:
            sp = Path(chosen)
            if not sp.is_absolute():
                root_rel = Path(cfg["_project_root"]) / sp
                sp = root_rel if root_rel.exists() else fig_dir / sp
            if not sp.exists():
                report["errors"].append(f"summary {pid}: scheme_image 不存在 {sp}")
            else:
                info = figs_info.get(pid) or {}
                if info.get("confidence") == "low" and not s.get("reviewed"):
                    report["errors"].append(f"summary {pid}: 低置信 Scheme 未经视觉确认")
        elif s.get("pdf_status") == "pdf_read":
            default = fig_dir / f"{pid}_fig1.png"
            if not default.exists():
                report["errors"].append(f"summary {pid}: pdf_read 但缺少 Scheme 图")

    # 4) docx
    try:
        from docx import Document
        d = Document(str(docx))
        report["docx"] = {
            "exists": True, "path": str(docx),
            "paragraphs": len(d.paragraphs), "tables": len(d.tables),
            "inline_shapes": len(d.inline_shapes),
        }
    except Exception as e:
        report["docx"] = {"exists": False, "path": str(docx), "error": str(e)}
        report["errors"].append(f"docx: {e}")

    # 5) pptx + 每篇2页结构抽查
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx))
        pics = sum(1 for sl in prs.slides for sh in sl.shapes if sh.shape_type == 13)
        slides = len(prs.slides._sldIdLst)
        n_summaries = len(report["summaries"])
        expected = REQUIRED_SLIDE_COUNT + 2 * n_summaries
        structure_ok = slides == expected
        report["pptx"] = {"exists": True, "path": str(pptx), "slides": slides,
                          "pictures": pics, "expected_slides": expected,
                          "per_paper_2_pages": structure_ok}
        if not structure_ok:
            report["errors"].append(
                f"pptx slides={slides}, expected={expected} (8 + 2×{n_summaries})")
    except Exception as e:
        report["pptx"] = {"exists": False, "path": str(pptx), "error": str(e)}
        report["errors"].append(f"pptx: {e}")

    report["ok"] = len(report["errors"]) == 0
    out = project_path(cfg, "check_report.json")
    write_json(out, report)
    print("check_report ->", out)
    for e in report["errors"]:
        print("  [FAIL]", e)
    print("OK" if report["ok"] else "HAS ERRORS")


if __name__ == "__main__":
    main()
