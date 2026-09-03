# -*- coding: utf-8 -*-
"""
check_layout.py — 用 python-pptx 检查每页文本是否可能溢出文本框/幻灯片边界（近似字符宽度估算）。
输出 layout_report.json：每个违约的 shape 记录。
PPTX 文件名由 project.yml / papers_meta 驱动。
"""
import math

from pptx import Presentation

from prk_config import load_papers_meta, output_dir, parse_project_arg, project_path, project_title, write_json


def est_overflow(slide_no, shape):
    tf = shape.text_frame
    w_in = shape.width / 914400.0 if shape.width else 1
    h_in = shape.height / 914400.0 if shape.height else 1
    total_lines = 0
    for para in tf.paragraphs:
        text = "".join(r.text for r in para.runs)
        if not text:
            total_lines += 0.3
            continue
        size = 12.0
        for r in para.runs:
            if r.font.size:
                size = r.font.size.pt
                break
        cpl = max(1, int(w_in * 72 / (size * 1.02)))
        lines = max(1, math.ceil(len(text) / cpl))
        total_lines += lines * (size * 1.35) / 72.0
    est_h = total_lines
    return est_h > h_in * 1.15 + 0.05, round(est_h, 2), round(h_in, 2)


def main():
    cfg, args = parse_project_arg()
    meta = load_papers_meta(cfg)
    title = project_title(cfg, meta)
    pptx = output_dir(cfg, "deliverables") / f"{title}_研究报告.pptx"
    if not pptx.exists():
        raise SystemExit(f"PPTX 不存在：{pptx}")

    prs = Presentation(str(pptx))
    issues = []
    total_shapes = 0
    for i, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            total_shapes += 1
            try:
                overflow, est, box = est_overflow(i, sh)
            except Exception:
                continue
            if overflow and est > box * 1.35:
                issues.append({"slide": i, "shape": sh.name,
                               "est_in": est, "box_in": box,
                               "text": sh.text_frame.text[:60]})
    report = {"slides": len(prs.slides._sldIdLst), "total_text_shapes": total_shapes,
              "probable_overflows": len(issues), "issues": issues[:20]}
    out = project_path(cfg, "layout_report.json")
    write_json(out, report)
    print("layout_report ->", out)


if __name__ == "__main__":
    main()
